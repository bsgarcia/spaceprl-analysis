from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import matplotlib.pyplot as plt

FOLDER = '../outputs/ppt/'

def add(filename, tight_layout=True):
   
    # if the file exists, open it, otherwise create a new one
    filename = FOLDER + filename
    try:
        prs = Presentation(filename)
    except:
        prs = Presentation()

    # select blank slide layout
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(0.5)
    top = Inches(0.5)
    
    # Calculate the width and height of the slide
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Calculate the width and height of the picture
    img_width, img_height = plt.gcf().get_size_inches()

    # Calculate the scaling factors
    scale_width = slide_width / Inches(img_width)
    scale_height = slide_height / Inches(img_height)

    # Choose the smaller scaling factor to fit the picture within the slide
    scale_factor = min(scale_width, scale_height)

    # Calculate the new width and height of the picture
    new_width = img_width * scale_factor
    new_height = img_height * scale_factor
    new_height*=0.9
    new_width*=0.9

    # Add the plot image from BytesIO to the slide and resize it
    left = (slide_width - Inches(new_width)) / 2
    top = (slide_height - Inches(new_height)) / 2

    buffer = BytesIO()
    # save and remove white space
    if tight_layout:
        plt.tight_layout()
    plt.savefig(buffer, format='png')
    buffer.seek(0)
    
    slide.shapes.add_picture(buffer, left, top, width=Inches(new_width), height=Inches(new_height))
    
    prs.save(filename)


